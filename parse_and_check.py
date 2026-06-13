"""把附件（xlsx / docx / pdf / pptx / html / md）解析成文本，并可选地跑一组校验脚本。

用法：
    python3 parse_and_check.py --file <path> [--script <path>]... [--out <path>] [--report <path>]

退出码：
    0  解析成功且每项 check 都通过（或者根本没传 --script）
    1  解析成功，但至少一项 check 返回 ok=False 或抛异常
    2  主脚本本身出错（文件不存在、不支持的后缀、parser 异常、--script 无法加载等）
"""

from __future__ import annotations

import argparse
import importlib.util
import json
import sys
import traceback
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Callable


# ---------- parsers ----------

def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
    out: list[str] = []
    for sheet in wb.worksheets:
        out.append(f"## sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append("\t".join(cells))
        out.append("")
    return "\n".join(out).rstrip() + "\n"


def _parse_docx(path: Path) -> str:
    import docx  # python-docx

    doc = docx.Document(str(path))
    out: list[str] = []

    # Walk body in document order so paragraphs and tables interleave correctly.
    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}", 1)[-1]
        if tag == "p":
            # Find matching paragraph object — python-docx exposes them on doc.paragraphs
            # but in body order. Cheap re-lookup via _element identity.
            for p in doc.paragraphs:
                if p._element is child:
                    text = p.text.rstrip()
                    if text:
                        style = (p.style.name or "").lower() if p.style else ""
                        if style.startswith("heading"):
                            try:
                                level = int(style.split()[-1])
                            except (ValueError, IndexError):
                                level = 1
                            out.append("#" * max(1, min(level, 6)) + " " + text)
                        else:
                            out.append(text)
                    break
        elif tag == "tbl":
            for tbl in doc.tables:
                if tbl._element is child:
                    rows = [
                        [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        for row in tbl.rows
                    ]
                    if rows:
                        widths = [max(len(r[i]) for r in rows) for i in range(len(rows[0]))]
                        for i, row in enumerate(rows):
                            out.append("| " + " | ".join(c.ljust(widths[j]) for j, c in enumerate(row)) + " |")
                            if i == 0:
                                out.append("|" + "|".join("-" * (w + 2) for w in widths) + "|")
                        out.append("")
                    break
    return "\n".join(out).rstrip() + "\n"


def _parse_pdf(path: Path) -> str:
    import pdfplumber

    out: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            out.append(f"--- page {i} ---")
            text = page.extract_text() or ""
            out.append(text.rstrip())
            out.append("")
    return "\n".join(out).rstrip() + "\n"


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text_frame.text or "").strip()
        out.append(f"## slide {i}" + (f": {title}" if title else ""))
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append("- " + text)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    out.append("| " + " | ".join(cells) + " |")
        # Notes
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                out.append("> notes: " + notes.replace("\n", " "))
        out.append("")
    return "\n".join(out).rstrip() + "\n"


def _parse_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    raw = path.read_bytes()
    soup = BeautifulSoup(raw, "html.parser")
    for bad in soup(["script", "style", "noscript"]):
        bad.decompose()

    # Promote semantic headings to markdown ATX so downstream text keeps the signal.
    for level in range(1, 7):
        for h in soup.find_all(f"h{level}"):
            text = h.get_text(" ", strip=True)
            h.replace_with(f"\n\n{'#' * level} {text}\n\n")

    # Title element gets pulled in too — many layout-only pages put the only real
    # heading in <title>, and we already see <body> headings via the H1-6 pass.
    title_el = soup.find("title")
    title_text = title_el.get_text(" ", strip=True) if title_el else ""

    text = soup.get_text("\n")
    if title_text:
        text = f"# {title_text}\n\n" + text

    # Collapse runs of blank lines.
    lines = [ln.rstrip() for ln in text.splitlines()]
    cleaned: list[str] = []
    blank = False
    for ln in lines:
        if ln.strip():
            cleaned.append(ln)
            blank = False
        elif not blank:
            cleaned.append("")
            blank = True
    return "\n".join(cleaned).strip() + "\n"


def _parse_md(path: Path) -> str:
    return path.read_text(encoding="utf-8")


PARSERS: dict[str, Callable[[Path], str]] = {
    ".xlsx": _parse_xlsx,
    ".docx": _parse_docx,
    ".pdf": _parse_pdf,
    ".pptx": _parse_pptx,
    ".html": _parse_html,
    ".htm": _parse_html,
    ".md": _parse_md,
    ".markdown": _parse_md,
}


# ---------- check runner ----------

@dataclass
class CheckResult:
    name: str
    status: str  # "ok" | "fail" | "error"
    detail: str
    script: str

    @property
    def label(self) -> str:
        return {"ok": "OK", "fail": "FAIL", "error": "ERR"}[self.status]


def _load_check(script_path: Path):
    # Make `from checks.checker import ContentChecker` resolvable in scripts no
    # matter the caller's cwd: project root is the dir holding this file.
    project_root = str(Path(__file__).resolve().parent)
    if project_root not in sys.path:
        sys.path.insert(0, project_root)

    spec = importlib.util.spec_from_file_location(f"_check_{script_path.stem}", str(script_path))
    if spec is None or spec.loader is None:
        raise ImportError(f"could not build import spec for {script_path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)

    # Prefer a Checker subclass defined in this module. We filter to classes
    # whose __module__ matches the loaded module so that re-imports of the
    # base classes (e.g. `from checks.checker import ContentChecker`) don't
    # accidentally show up here.
    try:
        from checks.checker import Checker  # runtime import: optional dep on the package
    except Exception:
        Checker = None  # type: ignore[assignment]

    checker_cls = None
    if Checker is not None:
        candidates = [
            obj
            for _, obj in vars(mod).items()
            if isinstance(obj, type)
            and issubclass(obj, Checker)
            and obj is not Checker
            and obj.__module__ == mod.__name__
            and not getattr(obj, "__abstract__", False)
        ]
        if len(candidates) > 1:
            raise RuntimeError(
                f"{script_path} defines multiple Checker subclasses: "
                + ", ".join(c.__name__ for c in candidates)
                + " — split them into separate scripts"
            )
        if candidates:
            checker_cls = candidates[0]

    if checker_cls is not None:
        return checker_cls()  # callable: instance(text, meta) -> dict

    fn = getattr(mod, "check", None)
    if not callable(fn):
        raise AttributeError(
            f"{script_path} must define a Checker subclass or a `check(text, meta)` function"
        )
    return fn


def _run_check(script_path: Path, text: str, meta: dict) -> CheckResult:
    default_name = script_path.stem
    try:
        fn = _load_check(script_path)
    except Exception as exc:
        return CheckResult(default_name, "error", f"load failed: {exc}", str(script_path))

    # If the loader returned a Checker instance, prefer its `.name` over the
    # script filename as the fallback label.
    instance_name = getattr(fn, "name", None)
    if isinstance(instance_name, str) and instance_name:
        default_name = instance_name

    try:
        raw = fn(text, meta)
    except Exception:
        tb = traceback.format_exc().strip().splitlines()
        return CheckResult(default_name, "error", "raised: " + " | ".join(tb[-2:]), str(script_path))

    if not isinstance(raw, dict) or "ok" not in raw:
        return CheckResult(default_name, "error", f"bad return shape: {raw!r}", str(script_path))

    name = raw.get("name") or default_name
    status = "ok" if raw["ok"] else "fail"
    detail = str(raw.get("detail", ""))
    return CheckResult(name, status, detail, str(script_path))


# ---------- main ----------

def _build_meta(path: Path) -> dict:
    return {
        "file_path": str(path),
        "file_name": path.name,
        "ext": path.suffix.lower(),
        "size_bytes": path.stat().st_size,
    }


def _default_out(file_path: Path) -> Path:
    return file_path.with_suffix(file_path.suffix + ".md") if file_path.suffix.lower() in (".md", ".markdown") else file_path.with_suffix(".md")


def _default_report(file_path: Path) -> Path:
    return Path("check_result") / _report_stem(file_path)


def _report_stem(file_path: Path) -> str:
    # 完整 file_name（含后缀）+ 时间戳前缀，避免跨格式同名互相覆盖。
    # 例：foo.docx → 20260613-155500__foo.docx.report.json
    from datetime import datetime
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    return f"{ts}__{file_path.name}.report.json"


def _resolve_report_path(arg: Path | None, file_path: Path) -> Path:
    """``--report`` 既可以传目录，也可以传具体的 .json 文件路径。

    - 不传（None）：用 ``check_result/<timestamp>__<file_name>.report.json``
    - 传的路径以 ``.json`` 结尾：当作完整文件路径，**不**自动加时间戳，
      如果已存在就在 stem 末尾追加 ``-2`` / ``-3`` …直到独占
    - 否则：当作目录，里面放 ``<timestamp>__<file_name>.report.json``

    无论哪种情况，最终路径都保证不会覆盖已有文件。
    """
    if arg is None:
        candidate = _default_report(file_path)
    elif arg.suffix.lower() == ".json":
        candidate = arg
    else:
        candidate = arg / _report_stem(file_path)
    return _avoid_overwrite(candidate)


def _avoid_overwrite(path: Path) -> Path:
    if not path.exists():
        return path
    # foo.report.json → foo.report-2.json / foo.report-3.json / ...
    # foo.json        → foo-2.json / foo-3.json / ...
    parent = path.parent
    name = path.name
    if name.endswith(".report.json"):
        base = name[: -len(".report.json")]
        suffix = ".report.json"
    else:
        base = path.stem
        suffix = path.suffix
    n = 2
    while True:
        candidate = parent / f"{base}-{n}{suffix}"
        if not candidate.exists():
            return candidate
        n += 1


def _expand_scripts(scripts: list[Path]) -> tuple[list[Path], list[str], list[str]]:
    """展开 ``--script`` 入参：

    - 是文件 → 原样保留
    - 是目录 → 顶层 ``check_*.py`` 升序展开，不递归
    - 不存在 / 是目录但里面一个 check 都没有 → 收集成错误信息

    返回 ``(expanded, warnings, errors)``。``warnings`` 用于"目录里有 ``.py`` 但
    命名不符合 ``check_*.py`` 规范"等不致命情况——主脚本会原样打到 stderr 让人看到，
    但不影响退出码。``errors`` 走 stderr + exit 2。
    """
    expanded: list[Path] = []
    warnings: list[str] = []
    errors: list[str] = []
    for entry in scripts:
        if not entry.exists():
            errors.append(f"--script not found: {entry}")
            continue
        if entry.is_file():
            expanded.append(entry)
            continue
        if entry.is_dir():
            top_py = sorted(p for p in entry.glob("*.py") if p.is_file())
            children = [p for p in top_py if p.name.startswith("check_")]
            skipped = [p for p in top_py if not p.name.startswith("check_")]
            if not children:
                errors.append(
                    f"--script directory has no check_*.py: {entry}"
                    + (f" (skipped non-conforming: {[p.name for p in skipped]})" if skipped else "")
                )
                continue
            for p in skipped:
                warnings.append(
                    f"skipping {p} — only check_*.py at the top level is auto-loaded"
                )
            expanded.extend(children)
            continue
        errors.append(f"--script is neither file nor directory: {entry}")
    return expanded, warnings, errors


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Parse an artifact to text and optionally run check scripts.")
    parser.add_argument("--file", required=True, type=Path, help="Path to artifact file (xlsx/docx/pdf/pptx/html/md).")
    parser.add_argument(
        "--script",
        action="append",
        default=[],
        type=Path,
        help=(
            "校验脚本路径，可重复传。可以是单个 .py 文件，也可以是一个目录"
            "（目录里顶层匹配 check_*.py 的脚本会按文件名升序依次执行，不递归）。"
        ),
    )
    parser.add_argument("--out", type=Path, default=None, help="Where to write parsed text. Default: <file>.md next to the input.")
    parser.add_argument(
        "--report",
        type=Path,
        default=None,
        help=(
            "JSON 报告输出位置；可传目录（落到 <dir>/<file_stem>.report.json）"
            "或直接传 .json 文件路径。默认 check_result/<file_stem>.report.json。"
            "只在传了 --script 时生成。"
        ),
    )
    args = parser.parse_args(argv)

    file_path: Path = args.file
    if not file_path.exists():
        print(f"error: file not found: {file_path}", file=sys.stderr)
        return 2
    if not file_path.is_file():
        print(f"error: not a regular file: {file_path}", file=sys.stderr)
        return 2

    ext = file_path.suffix.lower()
    parser_fn = PARSERS.get(ext)
    if parser_fn is None:
        print(f"error: unsupported extension {ext!r}; supported: {sorted(PARSERS)}", file=sys.stderr)
        return 2

    try:
        text = parser_fn(file_path)
    except Exception:
        print(f"error: parser failed for {file_path}", file=sys.stderr)
        traceback.print_exc()
        return 2

    out_path: Path = args.out if args.out is not None else _default_out(file_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(text, encoding="utf-8")

    meta = _build_meta(file_path)

    print(f"parsed {file_path} ({meta['size_bytes']} bytes, {len(text)} chars) → {out_path}")

    scripts: list[Path] = list(args.script or [])
    if not scripts:
        return 0

    # 把目录展开成具体的 check_*.py，并把"文件不存在 / 空目录"等错误一次报齐。
    expanded, warnings, errors = _expand_scripts(scripts)
    for msg in warnings:
        print(f"warning: {msg}", file=sys.stderr)
    if errors:
        for msg in errors:
            print(f"error: {msg}", file=sys.stderr)
        return 2

    results = [_run_check(p, text, meta) for p in expanded]

    n_ok = sum(1 for r in results if r.status == "ok")
    n_fail = sum(1 for r in results if r.status == "fail")
    n_err = sum(1 for r in results if r.status == "error")
    for r in results:
        line = f"[{r.label}] {r.name}"
        if r.detail:
            line += f" — {r.detail}"
        print(line)
    print(f"summary: {n_ok} pass / {n_fail} fail / {n_err} err")

    report_path: Path = _resolve_report_path(args.report, file_path)
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(
        json.dumps(
            {
                "file": meta,
                "out": str(out_path),
                "summary": {"pass": n_ok, "fail": n_fail, "err": n_err},
                "checks": [asdict(r) for r in results],
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )

    return 0 if (n_fail == 0 and n_err == 0) else 1


if __name__ == "__main__":
    sys.exit(main())
